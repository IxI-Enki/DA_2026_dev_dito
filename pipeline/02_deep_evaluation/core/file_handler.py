"""
File Handler - Robuster Dateizugriff

Liest verschiedene Dateiformate sicher ein und extrahiert Text/Metadaten.
Alle Einstellungen werden aus config/env.yaml geladen - KEINE hardcoded Werte!
"""

import json
import logging
from pathlib import Path
from typing import Dict, Any, Optional, Tuple, List, Set

# Optional dependencies
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import docx
except ImportError:
    docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

logger = logging.getLogger(__name__)


class FileHandler:
    """
    Handled das Lesen verschiedener Dateitypen.
    
    Alle Einstellungen werden aus config/env.yaml geladen.
    """

    def __init__(self, config: Optional[Any] = None):
        """
        Initialisiert den FileHandler.
        
        Args:
            config: EvaluationConfig Instanz (optional, wird bei Bedarf geladen)
        """
        # Lazy import um circular dependencies zu vermeiden
        if config is None:
            import sys
            sys.path.insert(0, str(Path(__file__).parent.parent))
            from config import get_config
            config = get_config()
        
        self.config = config
        self.raw_config = config.raw_config
        
        # Load encodings from config or use defaults
        format_cfg = self.raw_config.get('FORMAT_ANALYSIS', {})
        self.encodings = format_cfg.get('text_encodings', ['utf-8', 'latin-1', 'cp1252'])
        
        # Load extraction limits from config
        extraction_cfg = format_cfg.get('extraction_limits', {})
        self.max_pdf_pages = extraction_cfg.get('pdf_max_pages', 5)
        self.max_xlsx_rows = extraction_cfg.get('xlsx_max_rows', 50)
        
        # Load supported formats from config
        supported_formats = format_cfg.get('supported_formats', {})
        
        # Text formats
        text_exts: Set[str] = set()
        if 'pages' in supported_formats:
            page_exts = supported_formats['pages'].get('extensions', ['.txt', '.md'])
            if isinstance(page_exts, list):
                text_exts.update(page_exts)
        
        # Add common text formats
        text_exts.update(['.json', '.xml', '.html', '.css', '.js', '.php', '.yaml', '.yml'])
        
        self.text_extensions = text_exts
        
        # Document extensions
        doc_exts: Set[str] = set()
        if 'documents' in supported_formats:
            docs_cfg = supported_formats['documents']
            if 'pdf' in docs_cfg:
                pdf_exts = docs_cfg['pdf'].get('extensions', ['.pdf'])
                if isinstance(pdf_exts, list):
                    doc_exts.update(pdf_exts)
            if 'office' in docs_cfg:
                for office_type in docs_cfg['office'].values():
                    office_exts = office_type.get('extensions', [])
                    if isinstance(office_exts, list):
                        doc_exts.update(office_exts)
        
        self.document_extensions = doc_exts
        
        # Image extensions
        img_exts: Set[str] = set()
        if 'images' in supported_formats:
            img_exts_list = supported_formats['images'].get('extensions', ['.jpg', '.jpeg', '.png', '.svg'])
            if isinstance(img_exts_list, list):
                img_exts.update(img_exts_list)
        
        self.image_extensions = img_exts

    def read_text(self, file_path: Path) -> str:
        """
        Liest Textdatei mit Encoding-Fallback.
        
        Args:
            file_path: Pfad zur Textdatei
            
        Returns:
            Dateiinhalt als String, leerer String bei Fehler
        """
        for encoding in self.encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.error(f"Error reading {file_path}: {e}")
                return ""
        
        logger.error(f"Failed to decode {file_path} with encodings: {self.encodings}")
        return ""

    def read_json(self, file_path: Path) -> Dict[str, Any]:
        """
        Liest JSON sicher ein.
        
        Args:
            file_path: Pfad zur JSON-Datei
            
        Returns:
            Dictionary mit JSON-Inhalt, leeres Dict bei Fehler
        """
        text = self.read_text(file_path)
        if not text:
            return {}
        try:
            return json.loads(text)
        except json.JSONDecodeError as e:
            logger.error(f"JSON decode error in {file_path}: {e}")
            return {}

    def extract_text_from_pdf(self, file_path: Path, max_pages: Optional[int] = None) -> str:
        """
        Extrahiert Text aus PDF (benötigt PyMuPDF).
        
        Args:
            file_path: Pfad zur PDF-Datei
            max_pages: Maximale Anzahl Seiten (None = aus config)
            
        Returns:
            Extrahierter Text oder Fehlermeldung
        """
        if not fitz:
            return "[Error: PyMuPDF not installed]"
        
        if max_pages is None:
            max_pages = self.max_pdf_pages
        
        text = ""
        try:
            doc = fitz.open(file_path)
            try:
                total_pages = len(doc)
                pages_to_read = min(max_pages, total_pages) if max_pages is not None else total_pages
                for i in range(pages_to_read):
                    page = doc[i]
                    page_text = page.get_text()
                    if isinstance(page_text, str):
                        text += page_text + "\n"
                if max_pages is not None and total_pages > max_pages:
                    text += f"\n... (truncated after {max_pages}/{total_pages} pages)"
            finally:
                doc.close()
        except Exception as e:
            logger.error(f"Error extracting PDF {file_path}: {e}")
            return f"[Error extracting PDF: {e}]"
        
        return text

    def extract_text_from_docx(self, file_path: Path) -> str:
        """
        Extrahiert Text aus DOCX (benötigt python-docx).
        
        Args:
            file_path: Pfad zur DOCX-Datei
            
        Returns:
            Extrahierter Text oder Fehlermeldung
        """
        if not docx:
            return "[Error: python-docx not installed]"
        
        try:
            doc = docx.Document(str(file_path))
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            logger.error(f"Error extracting DOCX {file_path}: {e}")
            return f"[Error extracting DOCX: {e}]"

    def extract_text_from_xlsx(self, file_path: Path, max_rows: Optional[int] = None) -> str:
        """
        Extrahiert Text aus XLSX (benötigt openpyxl).
        
        Args:
            file_path: Pfad zur XLSX-Datei
            max_rows: Maximale Anzahl Zeilen pro Sheet (None = aus config)
            
        Returns:
            Extrahierter Text oder Fehlermeldung
        """
        if not openpyxl:
            return "[Error: openpyxl not installed]"
        
        if max_rows is None:
            max_rows = self.max_xlsx_rows
        
        text = ""
        try:
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            for sheet in wb.worksheets:
                text += f"Sheet: {sheet.title}\n"
                for i, row in enumerate(sheet.rows):
                    if max_rows is not None and i >= max_rows:
                        text += f"... (truncated after {max_rows} rows)\n"
                        break
                    row_text = " | ".join([str(cell.value) for cell in row if cell.value])
                    if row_text:
                        text += row_text + "\n"
        except Exception as e:
            logger.error(f"Error extracting XLSX {file_path}: {e}")
            return f"[Error extracting XLSX: {e}]"
        
        return text

    def extract_text_from_pptx(self, file_path: Path) -> str:
        """
        Extrahiert Text aus PPTX (benötigt python-pptx).
        
        Args:
            file_path: Pfad zur PPTX-Datei
            
        Returns:
            Extrahierter Text oder Fehlermeldung
        """
        if not Presentation:
            return "[Error: python-pptx not installed]"
        
        text = ""
        try:
            prs = Presentation(str(file_path))
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    # Check if shape has text attribute (BaseShape doesn't always have it)
                    # Type ignore because python-pptx BaseShape type hints are incomplete
                    try:
                        shape_text = getattr(shape, "text", None)  # type: ignore
                        if shape_text:
                            text += str(shape_text) + "\n"
                    except (AttributeError, TypeError):
                        pass
        except Exception as e:
            logger.error(f"Error extracting PPTX {file_path}: {e}")
            return f"[Error extracting PPTX: {e}]"
        
        return text

    def get_file_content(self, file_path: Path) -> str:
        """
        Generischer Wrapper für Content-Extraktion basierend auf Extension.
        
        Args:
            file_path: Pfad zur Datei
            
        Returns:
            Extrahierter Text oder Fehlermeldung
        """
        suffix = file_path.suffix.lower()
        
        if suffix == '.pdf':
            return self.extract_text_from_pdf(file_path)
        elif suffix == '.docx':
            return self.extract_text_from_docx(file_path)
        elif suffix == '.xlsx':
            return self.extract_text_from_xlsx(file_path)
        elif suffix == '.pptx':
            return self.extract_text_from_pptx(file_path)
        elif suffix in self.text_extensions:
            return self.read_text(file_path)
        else:
            return f"[Unsupported format: {suffix}]"

    def is_supported_format(self, file_path: Path) -> bool:
        """
        Prüft ob ein Dateiformat unterstützt wird.
        
        Args:
            file_path: Pfad zur Datei
            
        Returns:
            True wenn Format unterstützt wird
        """
        suffix = file_path.suffix.lower()
        all_supported = self.text_extensions | self.document_extensions | self.image_extensions
        return suffix in all_supported


# Test
if __name__ == "__main__":
    # Check imports
    print("=" * 60)
    print("  FILE HANDLER DEPENDENCIES")
    print("=" * 60)
    print(f"PyMuPDF available:     {fitz is not None}")
    print(f"python-docx available: {docx is not None}")
    print(f"openpyxl available:    {openpyxl is not None}")
    print(f"python-pptx available: {Presentation is not None}")
    print("=" * 60)
    
    # Test config loading
    try:
        handler = FileHandler()
        print(f"\n  Encodings: {handler.encodings}")
        print(f"  Max PDF pages: {handler.max_pdf_pages}")
        print(f"  Max XLSX rows: {handler.max_xlsx_rows}")
        print(f"  Text extensions: {sorted(handler.text_extensions)}")
        print(f"  Document extensions: {sorted(handler.document_extensions)}")
        print(f"  Image extensions: {sorted(handler.image_extensions)}")
    except Exception as e:
        print(f"\n  Error loading config: {e}")
