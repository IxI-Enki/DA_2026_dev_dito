#!/usr/bin/env python3
"""
RAG Preprocessing Pipeline - Main Entry Point
==============================================

Converts fetched DokuWiki content to RAG-optimized Markdown with YAML frontmatter.

Usage:
    python main.py                          # Auto-detect latest fetch
    python main.py --input-dir <path>       # Specify input directory
    python main.py --help                   # Show help

Input:  data/fetched/fetched_at_*/
Output: data/preprocessed/preprocess_at_*/
"""

import os
import sys
import json
import logging
import argparse
from pathlib import Path
from datetime import datetime
from typing import Dict, Any, Optional, List

# Add parent to path for imports
sys.path.insert(0, str(Path(__file__).parent))

from config import get_config, get_latest_fetch_dir, get_latest_evaluation
from page_processor import PageProcessor
from metadata_enricher import MetadataEnricher, MediaMetadataEnricher

logger = logging.getLogger(__name__)


class RAGPreprocessor:
    """Main orchestrator for RAG preprocessing pipeline."""
    
    def __init__(self, config_path: Optional[Path] = None):
        self.config = get_config(config_path)
        self.page_processor = PageProcessor(self.config.wiki_base_url)
        self.metadata_enricher = MetadataEnricher(self.config.wiki_base_url)
        self.media_enricher = MediaMetadataEnricher(self.config.wiki_base_url)
        
        # Statistics
        self.stats = {
            'pages_processed': 0,
            'pages_success': 0,
            'pages_failed': 0,
            'media_processed': 0,
            'media_success': 0,
            'media_failed': 0,
            'errors': [],
        }
    
    def run(
        self,
        input_dir: Optional[Path] = None,
        evaluation_file: Optional[Path] = None,
        output_dir: Optional[Path] = None,
    ) -> Dict[str, Any]:
        """
        Run the preprocessing pipeline.
        
        Args:
            input_dir: Fetched data directory (auto-detect if None)
            evaluation_file: Evaluation JSON from Stage 2 (auto-detect if None)
            output_dir: Output directory (auto-generate if None)
            
        Returns:
            Statistics dict
        """
        # Resolve input directory
        if input_dir is None:
            input_dir = get_latest_fetch_dir(self.config.fetched_dir)
            if input_dir is None:
                raise ValueError(f"No fetched data found in {self.config.fetched_dir}")
        
        logger.info(f"Input directory: {input_dir}")
        
        # Resolve evaluation file
        if evaluation_file is None:
            evaluation_file = get_latest_evaluation(self.config.evaluated_dir)
        
        if evaluation_file:
            logger.info(f"Loading evaluation: {evaluation_file}")
            self.metadata_enricher.load_evaluation(evaluation_file)
        else:
            logger.warning("No evaluation file found - proceeding without quality data")
        
        # Create output directory
        if output_dir is None:
            timestamp = datetime.now().strftime(self.config.output.get('timestamp_format', '%Y%m%d_%H%M%S'))
            prefix = self.config.output.get('prefix', 'preprocess_at_')
            output_dir = self.config.output_dir / f"{prefix}{timestamp}"
        
        output_dir.mkdir(parents=True, exist_ok=True)
        pages_dir = output_dir / self.config.output.get('pages_subdir', 'pages')
        media_dir = output_dir / self.config.output.get('media_subdir', 'media')
        pages_dir.mkdir(exist_ok=True)
        media_dir.mkdir(exist_ok=True)
        
        logger.info(f"Output directory: {output_dir}")
        
        # Extract fetch timestamp from input directory name
        fetch_timestamp = self._extract_timestamp(input_dir.name)
        
        # Process pages
        logger.info("Processing pages...")
        self._process_pages(input_dir, pages_dir, fetch_timestamp)
        
        # Process media (text extraction)
        logger.info("Processing media...")
        self._process_media(input_dir, media_dir, fetch_timestamp)
        
        # Generate manifest
        self._generate_manifest(output_dir, input_dir)
        
        # Print summary
        self._print_summary()
        
        return self.stats
    
    def _extract_timestamp(self, dir_name: str) -> Optional[str]:
        """Extract timestamp from directory name like fetched_at_20260201_120240."""
        if 'fetched_at_' in dir_name:
            ts_part = dir_name.replace('fetched_at_', '')
            try:
                dt = datetime.strptime(ts_part, '%Y%m%d_%H%M%S')
                return dt.isoformat()
            except ValueError:
                pass
        return None
    
    def _process_pages(self, input_dir: Path, output_dir: Path, fetch_timestamp: Optional[str]) -> None:
        """Process all pages from the fetched data."""
        page_content_dir = input_dir / 'page_content'
        raw_json_dir = input_dir / 'raw_json'
        page_links_dir = input_dir / 'page_links'
        
        if not page_content_dir.exists():
            logger.error(f"Page content directory not found: {page_content_dir}")
            return
        
        page_files = list(page_content_dir.glob('*.txt'))
        logger.info(f"Found {len(page_files)} page files")
        
        for page_file in page_files:
            self.stats['pages_processed'] += 1
            
            try:
                # Derive page_id from filename (e.g., 'namespace_page.txt' -> 'namespace:page')
                page_id = page_file.stem.replace('_', ':')
                
                # Read wiki content
                wiki_content = page_file.read_text(encoding='utf-8')
                
                # Skip empty pages
                if not wiki_content.strip():
                    logger.debug(f"Skipping empty page: {page_id}")
                    continue
                
                # Convert to Markdown
                result = self.page_processor.convert(wiki_content, page_id)
                
                if not result.success:
                    logger.warning(f"Conversion issues for {page_id}: {result.errors}")
                
                # Load metadata
                raw_metadata = self._load_json(raw_json_dir / f"{page_file.stem}_complete.json")
                links_data = self._load_json(page_links_dir / f"{page_file.stem}_links.json")
                
                # Generate frontmatter
                frontmatter = self.metadata_enricher.generate_frontmatter(
                    page_id=page_id,
                    title=result.title,
                    raw_metadata=raw_metadata,
                    links_data=links_data,
                    fetch_timestamp=fetch_timestamp,
                )
                
                # Combine frontmatter and content
                output_content = frontmatter + '\n' + result.markdown
                
                # Write output file
                output_file = output_dir / f"{page_file.stem}.md"
                output_file.write_text(output_content, encoding='utf-8')
                
                self.stats['pages_success'] += 1
                
            except Exception as e:
                logger.error(f"Failed to process {page_file.name}: {e}")
                self.stats['pages_failed'] += 1
                self.stats['errors'].append(f"Page {page_file.name}: {str(e)}")
                
                if len(self.stats['errors']) >= self.config.processing.get('max_errors', 50):
                    logger.error("Max errors reached, stopping")
                    break
    
    def _process_media(self, input_dir: Path, output_dir: Path, fetch_timestamp: Optional[str]) -> None:
        """Process media files - extract text from documents."""
        media_dir = input_dir / 'media'
        
        if not media_dir.exists():
            logger.warning(f"Media directory not found: {media_dir}")
            return
        
        # Find all media files recursively
        media_files = list(media_dir.rglob('*'))
        media_files = [f for f in media_files if f.is_file()]
        
        logger.info(f"Found {len(media_files)} media files")
        
        # Filter to extractable types
        extractable_extensions = {'.pdf', '.docx', '.xlsx', '.pptx', '.txt'}
        
        for media_file in media_files:
            if media_file.suffix.lower() not in extractable_extensions:
                continue
            
            self.stats['media_processed'] += 1
            
            try:
                # Derive media_id from relative path
                rel_path = media_file.relative_to(media_dir)
                media_id = str(rel_path).replace('\\', ':').replace('/', ':')
                media_id = media_id.rsplit('.', 1)[0]  # Remove extension
                
                # Extract text based on file type
                extracted_text = self._extract_media_text(media_file)
                
                if not extracted_text:
                    logger.debug(f"No text extracted from: {media_file.name}")
                    continue
                
                # Generate frontmatter
                frontmatter = self.media_enricher.generate_frontmatter(
                    media_id=media_id,
                    file_path=media_file,
                    file_size=media_file.stat().st_size,
                    fetch_timestamp=fetch_timestamp,
                )
                
                # Combine frontmatter and content
                output_content = frontmatter + '\n' + extracted_text
                
                # Write output file
                safe_name = media_file.stem.replace(' ', '_')
                output_file = output_dir / f"{safe_name}.txt"
                output_file.write_text(output_content, encoding='utf-8')
                
                self.stats['media_success'] += 1
                
            except Exception as e:
                logger.error(f"Failed to process media {media_file.name}: {e}")
                self.stats['media_failed'] += 1
    
    def _extract_media_text(self, file_path: Path) -> Optional[str]:
        """Extract text from a media file."""
        ext = file_path.suffix.lower()
        
        try:
            if ext == '.txt':
                return file_path.read_text(encoding='utf-8', errors='ignore')
            
            elif ext == '.pdf':
                return self._extract_pdf_text(file_path)
            
            elif ext == '.docx':
                return self._extract_docx_text(file_path)
            
            elif ext == '.xlsx':
                return self._extract_xlsx_text(file_path)
            
            elif ext == '.pptx':
                return self._extract_pptx_text(file_path)
            
        except ImportError as e:
            logger.warning(f"Missing library for {ext}: {e}")
        except Exception as e:
            logger.warning(f"Extraction failed for {file_path.name}: {e}")
        
        return None
    
    def _extract_pdf_text(self, file_path: Path) -> Optional[str]:
        """Extract text from PDF."""
        try:
            from pypdf import PdfReader
        except ImportError:
            from PyPDF2 import PdfReader
        
        reader = PdfReader(file_path)
        text_parts = []
        
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        
        return '\n\n'.join(text_parts) if text_parts else None
    
    def _extract_docx_text(self, file_path: Path) -> Optional[str]:
        """Extract text from DOCX."""
        from docx import Document
        
        doc = Document(str(file_path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        
        return '\n\n'.join(paragraphs) if paragraphs else None
    
    def _extract_xlsx_text(self, file_path: Path) -> Optional[str]:
        """Extract text from XLSX."""
        from openpyxl import load_workbook
        
        wb = load_workbook(file_path, read_only=True, data_only=True)
        text_parts = []
        
        for sheet in wb.worksheets:
            sheet_data = []
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else '' for cell in row]
                if any(row_text):
                    sheet_data.append(' | '.join(row_text))
            
            if sheet_data:
                text_parts.append(f"## {sheet.title}\n" + '\n'.join(sheet_data))
        
        return '\n\n'.join(text_parts) if text_parts else None
    
    def _extract_pptx_text(self, file_path: Path) -> Optional[str]:
        """Extract text from PPTX."""
        from pptx import Presentation
        
        prs = Presentation(str(file_path))
        text_parts = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text.strip():
                    slide_text.append(text)
            
            if slide_text:
                text_parts.append(f"## Slide {i}\n" + '\n'.join(slide_text))
        
        return '\n\n'.join(text_parts) if text_parts else None
    
    def _load_json(self, path: Path) -> Optional[Dict[str, Any]]:
        """Safely load JSON file."""
        if not path.exists():
            return None
        try:
            with open(path, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            logger.debug(f"Could not load {path}: {e}")
            return None
    
    def _generate_manifest(self, output_dir: Path, input_dir: Path) -> None:
        """Generate manifest.json with preprocessing metadata."""
        manifest = {
            'preprocessing_info': {
                'version': '1.0.0',
                'preprocessed_at': datetime.now().isoformat(),
                'input_dir': str(input_dir),
                'output_dir': str(output_dir),
            },
            'statistics': self.stats,
            'pages': [],
            'media': [],
        }
        
        # List generated files
        pages_dir = output_dir / self.config.output.get('pages_subdir', 'pages')
        media_dir = output_dir / self.config.output.get('media_subdir', 'media')
        
        if pages_dir.exists():
            manifest['pages'] = [f.name for f in pages_dir.glob('*.md')]
        
        if media_dir.exists():
            manifest['media'] = [f.name for f in media_dir.glob('*.txt')]
        
        # Write manifest
        manifest_file = output_dir / self.config.output.get('manifest_file', 'manifest.json')
        with open(manifest_file, 'w', encoding='utf-8') as f:
            json.dump(manifest, f, indent=2, ensure_ascii=False)
        
        logger.info(f"Manifest written: {manifest_file}")
    
    def _print_summary(self) -> None:
        """Print processing summary."""
        print("\n" + "=" * 60)
        print("RAG PREPROCESSING COMPLETE")
        print("=" * 60)
        print(f"Pages processed: {self.stats['pages_processed']}")
        print(f"  - Success: {self.stats['pages_success']}")
        print(f"  - Failed:  {self.stats['pages_failed']}")
        print(f"Media processed: {self.stats['media_processed']}")
        print(f"  - Success: {self.stats['media_success']}")
        print(f"  - Failed:  {self.stats['media_failed']}")
        
        if self.stats['errors']:
            print(f"\nErrors ({len(self.stats['errors'])}):")
            for error in self.stats['errors'][:10]:
                print(f"  - {error}")
            if len(self.stats['errors']) > 10:
                print(f"  ... and {len(self.stats['errors']) - 10} more")
        
        print("=" * 60)


def setup_logging(level: str = 'INFO') -> None:
    """Setup logging configuration."""
    logging.basicConfig(
        level=getattr(logging, level.upper()),
        format='%(asctime)s [%(levelname)s] %(name)s: %(message)s',
        handlers=[
            logging.StreamHandler(sys.stdout),
        ]
    )


def main():
    """Main entry point."""
    parser = argparse.ArgumentParser(
        description='RAG Preprocessing Pipeline - Convert Wiki to Markdown',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python main.py                              # Auto-detect latest fetch
  python main.py --input-dir data/fetched/fetched_at_20260201_120240
  python main.py --evaluation-file data/evaluated/evaluation_*.json
        """
    )
    
    parser.add_argument(
        '--input-dir', '-i',
        type=Path,
        help='Input directory (fetched_at_* folder)'
    )
    
    parser.add_argument(
        '--evaluation-file', '-e',
        type=Path,
        help='Evaluation JSON from Stage 2'
    )
    
    parser.add_argument(
        '--output-dir', '-o',
        type=Path,
        help='Output directory'
    )
    
    parser.add_argument(
        '--config', '-c',
        type=Path,
        help='Configuration file (env.yaml)'
    )
    
    parser.add_argument(
        '--log-level', '-l',
        default='INFO',
        choices=['DEBUG', 'INFO', 'WARNING', 'ERROR'],
        help='Logging level'
    )
    
    args = parser.parse_args()
    
    # Setup logging
    setup_logging(args.log_level)
    
    # Run preprocessing
    try:
        preprocessor = RAGPreprocessor(args.config)
        stats = preprocessor.run(
            input_dir=args.input_dir,
            evaluation_file=args.evaluation_file,
            output_dir=args.output_dir,
        )
        
        # Exit with error code if failures
        if stats['pages_failed'] > 0 or stats['media_failed'] > 0:
            sys.exit(1)
        
    except Exception as e:
        logger.error(f"Pipeline failed: {e}")
        raise


if __name__ == '__main__':
    main()
