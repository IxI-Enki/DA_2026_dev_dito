"""Media Processor (T071-T073, T033/US9)

Processes media files for the RAG pipeline:
- PDF text extraction with OCR fallback + quality cleaning (US7)
- DOCX, XLSX, PPTX text extraction (US9 -- migrated from main.py)
- Image OCR via Tesseract
- Batch processing of media directories
"""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# Supported file extensions
PDF_EXTENSIONS = {".pdf"}
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".gif"}
TEXT_EXTENSIONS = {".txt"}


class MediaProcessor:
    """Processes media files for RAG pipeline.

    Args:
        tesseract_path: Path to the Tesseract OCR binary.
            Defaults to system PATH lookup.
        ocr_language: Tesseract language string (e.g. ``deu+eng``).
    """

    def __init__(
        self,
        tesseract_path: str = "",
        ocr_language: str = "deu+eng",
    ) -> None:
        self.tesseract_path = tesseract_path
        self.ocr_language = ocr_language

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def process_pdf(self, pdf_path: Path) -> str:
        """Extract text from PDF. Docling first (layout-aware Markdown), then pypdf+cleanup, then OCR.

        Args:
            pdf_path: Path to the PDF file.

        Returns:
            Extracted text (empty string on failure).
        """
        if not pdf_path.exists():
            logger.warning("PDF not found: %s", pdf_path)
            return ""

        # 1. Try Docling first (layout-aware, structured Markdown; skip post-processing)
        text = self._extract_pdf_text_docling(pdf_path)
        if text.strip():
            return text

        # 2. Fallback: pypdf + cleanup heuristics
        text = self._extract_pdf_text(pdf_path)
        if text.strip():
            return self.clean_pdf_text(text)

        # 3. Fallback: OCR the PDF
        logger.info("PDF text empty, attempting OCR: %s", pdf_path.name)
        ocr_text = self._ocr_pdf(pdf_path)
        return self.clean_pdf_text(ocr_text) if ocr_text.strip() else ""

    def process_image(self, image_path: Path) -> str:
        """OCR an image using Tesseract.

        Args:
            image_path: Path to the image file.

        Returns:
            Extracted text (empty string on failure).
        """
        if not image_path.exists():
            logger.warning("Image not found: %s", image_path)
            return ""
        return self._ocr_image(image_path)

    def process_docx(self, docx_path: Path) -> str:
        """Extract text from a DOCX file. Docling first (structured Markdown), then python-docx with tables.

        Args:
            docx_path: Path to the DOCX file.

        Returns:
            Extracted text (empty string on failure).
        """
        if not docx_path.exists():
            logger.warning("DOCX not found: %s", docx_path)
            return ""
        # 1. Try Docling first (paragraphs + tables as structured Markdown)
        text = self._extract_docx_text_docling(docx_path)
        if text.strip():
            return text
        # 2. Fallback: python-docx with paragraphs and tables
        return self._extract_docx_text_fallback(docx_path)

    def process_xlsx(self, xlsx_path: Path) -> str:
        """Extract text from an XLSX file. Docling first (structured Markdown), then openpyxl with proper tables.

        Args:
            xlsx_path: Path to the XLSX file.

        Returns:
            Extracted text (empty string on failure).
        """
        if not xlsx_path.exists():
            logger.warning("XLSX not found: %s", xlsx_path)
            return ""
        # 1. Try Docling first (structured Markdown)
        text = self._extract_xlsx_text_docling(xlsx_path)
        if text.strip():
            return text
        # 2. Fallback: openpyxl with proper Markdown table syntax
        return self._extract_xlsx_text_fallback(xlsx_path)

    def process_pptx(self, pptx_path: Path) -> str:
        """Extract text from a PPTX file.

        Args:
            pptx_path: Path to the PPTX file.

        Returns:
            Extracted text (empty string on failure).
        """
        if not pptx_path.exists():
            logger.warning("PPTX not found: %s", pptx_path)
            return ""
        try:
            from pptx import Presentation

            prs = Presentation(str(pptx_path))
            text_parts: list[str] = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text: list[str] = []
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text.strip():
                        slide_text.append(text)
                if slide_text:
                    text_parts.append(f"## Slide {i}\n" + "\n".join(slide_text))
            return "\n\n".join(text_parts) if text_parts else ""
        except ImportError:
            logger.warning("python-pptx not installed -- PPTX extraction unavailable")
            return ""
        except Exception as e:
            logger.warning("PPTX extraction failed for %s: %s", pptx_path.name, e)
            return ""

    def process_media_directory(self, media_dir: Path) -> list[dict[str, Any]]:
        """Process all supported media files in a directory.

        Handles all document formats (PDF, DOCX, XLSX, PPTX) and images.

        Args:
            media_dir: Directory to scan.

        Returns:
            List of ``{filename, text, type}`` dicts.
        """
        results: list[dict[str, Any]] = []
        if not media_dir.exists():
            return results

        for f in sorted(media_dir.rglob("*")):
            if not f.is_file():
                continue
            ext = f.suffix.lower()
            text = ""
            ftype = ""
            if ext == ".pdf":
                text = self.process_pdf(f)
                ftype = "pdf"
            elif ext == ".docx":
                text = self.process_docx(f)
                ftype = "docx"
            elif ext == ".xlsx":
                text = self.process_xlsx(f)
                ftype = "xlsx"
            elif ext == ".pptx":
                text = self.process_pptx(f)
                ftype = "pptx"
            elif ext in IMAGE_EXTENSIONS:
                text = self.process_image(f)
                ftype = "image"
            else:
                continue

            results.append({"filename": f.name, "text": text, "type": ftype})
        return results

    # ------------------------------------------------------------------
    # PDF text quality (US7)
    # ------------------------------------------------------------------

    def clean_pdf_text(self, raw_text: str) -> str:
        """Post-process extracted PDF text for quality.

        Chains: fix initial char splits -> fix spaced characters -> merge short lines.
        """
        text = self._fix_initial_char_splits(raw_text)
        text = self._fix_spaced_characters(text)
        text = self._merge_short_lines(text)
        return text

    def _fix_initial_char_splits(self, text: str) -> str:
        """Rejoin PDF artifacts where the first 1-2 chars of a word are split off.

        Three patterns applied sequentially:

        A)  1-2 ALL-UPPERCASE letters + space + lowercase continuation (3+ chars).
            ``P rüfer`` -> ``Prüfer``, ``PR üfung`` -> ``PRüfung``.
        B)  1 uppercase + 1 lowercase (NOT a common German word) + space +
            lowercase continuation (5+ chars).
            ``Sc hriftlich`` -> ``Schriftlich``.
        C)  Single lowercase letter + space + umlaut/sharp-s start (3+ chars).
            ``m ündlich`` -> ``mündlich``.
        """
        if not text:
            return ""

        # Pattern A: 1-2 ALL uppercase + space + lowercase continuation (>= 3 chars)
        text = re.compile(
            r"\b([A-Z\u00C4\u00D6\u00DC]{1,2})\s+" r"([a-z\u00E4\u00F6\u00FC\u00DF]\w{2,})\b",
            re.UNICODE,
        ).sub(r"\1\2", text)

        # Pattern B: 1 uppercase + 1 lowercase + space + lowercase cont. (>= 5 chars)
        # Exclude common 2-letter German words (Da, Er, Es, Im, In, So, Um, Zu, ...)
        _COMMON_2 = frozenset(
            [
                "ab",
                "am",
                "an",
                "da",
                "du",
                "er",
                "es",
                "im",
                "in",
                "ja",
                "je",
                "na",
                "ob",
                "oh",
                "so",
                "um",
                "wo",
                "zu",
            ]
        )

        def _pattern_b_replace(m: re.Match[str]) -> str:
            frag = m.group(1)
            if frag.lower() in _COMMON_2:
                return m.group(0)  # keep original
            return frag + m.group(2)

        text = re.compile(
            r"\b([A-Z\u00C4\u00D6\u00DC][a-z\u00E4\u00F6\u00FC\u00DF])\s+"
            r"([a-z\u00E4\u00F6\u00FC\u00DF]\w{4,})\b",
            re.UNICODE,
        ).sub(_pattern_b_replace, text)

        # Pattern C: single lowercase letter + space + umlaut/sharp-s start (>= 3 chars)
        text = re.compile(
            r"\b([a-z])\s+([\u00E4\u00F6\u00FC\u00DF]\w{2,})\b",
            re.UNICODE,
        ).sub(r"\1\2", text)

        return text

    def _fix_spaced_characters(self, text: str) -> str:
        """Fix PDF layout artifacts with spaced characters.

        Heuristic: if >60% of 'words' on a line are single characters,
        join them and split on double-spaces.

        Example: ``"H T B L A  L e o n d i n g"`` -> ``"HTBLA Leonding"``
        """
        if not text:
            return ""

        lines = text.split("\n")
        fixed: list[str] = []

        for line in lines:
            words = line.split()
            if not words:
                fixed.append(line)
                continue

            single_chars = sum(1 for w in words if len(w) == 1)
            ratio = single_chars / len(words)

            if ratio > 0.6:
                # Split on double-space (word boundary in spaced text),
                # then join single chars within each group
                groups = re.split(r"  +", line.strip())
                rebuilt = []
                for group in groups:
                    parts = group.split()
                    if all(len(p) == 1 for p in parts):
                        rebuilt.append("".join(parts))
                    else:
                        rebuilt.append(group)
                fixed.append(" ".join(rebuilt))
            else:
                fixed.append(line)

        return "\n".join(fixed)

    def _merge_short_lines(self, text: str, threshold: int = 40) -> str:
        """Merge consecutive short lines into paragraphs.

        Respects structure: list items, headings, and empty lines are
        never merged with the previous line.
        """
        if not text:
            return ""

        lines = text.split("\n")
        merged: list[str] = []

        for line in lines:
            stripped = line.strip()

            # Empty lines are paragraph separators -- keep them
            if not stripped:
                merged.append("")
                continue

            # Structure-preserving: don't merge list items, headings
            is_structure = (
                stripped.startswith("#")
                or stripped.startswith("- ")
                or stripped.startswith("* ")
                or re.match(r"^\d+[\.\)]\s", stripped)
            )

            if is_structure:
                merged.append(stripped)
                continue

            # If previous line exists, is non-empty, and was short -> merge
            if (
                merged
                and merged[-1]
                and len(merged[-1]) < threshold
                and not merged[-1].startswith("#")
                and not merged[-1].startswith("- ")
                and not merged[-1].startswith("* ")
                and not re.match(r"^\d+[\.\)]\s", merged[-1])
            ):
                merged[-1] = merged[-1] + " " + stripped
            else:
                merged.append(stripped)

        return "\n".join(merged)

    # ------------------------------------------------------------------
    # Private helpers
    # ------------------------------------------------------------------

    def _extract_pdf_text_docling(self, pdf_path: Path) -> str:
        """Extract text from PDF using Docling (layout-aware, structured Markdown)."""
        try:
            from docling.document_converter import DocumentConverter

            converter = DocumentConverter()
            result = converter.convert(str(pdf_path))
            return result.document.export_to_markdown()
        except ImportError:
            logger.info("Docling not installed, falling back to pypdf")
            return ""
        except Exception as e:
            logger.warning("Docling failed for %s: %s, falling back", pdf_path.name, e)
            return ""

    def _extract_docx_text_docling(self, docx_path: Path) -> str:
        """Extract text from DOCX using Docling (paragraphs + tables as structured Markdown)."""
        try:
            from docling.document_converter import DocumentConverter

            converter = DocumentConverter()
            result = converter.convert(str(docx_path))
            return result.document.export_to_markdown()
        except ImportError:
            logger.info("Docling not installed for DOCX, falling back to python-docx")
            return ""
        except Exception as e:
            logger.warning("Docling failed for %s: %s, falling back", docx_path.name, e)
            return ""

    def _extract_docx_text_fallback(self, docx_path: Path) -> str:
        """Extract text from DOCX using python-docx (paragraphs + tables as Markdown)."""
        try:
            from docx import Document

            doc = Document(str(docx_path))
            parts: list[str] = []
            for p in doc.paragraphs:
                if p.text.strip():
                    parts.append(p.text.strip())
            for table in doc.tables:
                rows: list[str] = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                if rows:
                    ncols = len(table.columns)
                    header_sep = "| " + " | ".join(["---"] * ncols) + " |"
                    rows.insert(1, header_sep)
                    parts.append("\n".join(rows))
            return "\n\n".join(parts) if parts else ""
        except ImportError:
            logger.warning("python-docx not installed -- DOCX extraction unavailable")
            return ""
        except Exception as e:
            logger.warning("DOCX extraction failed for %s: %s", docx_path.name, e)
            return ""

    def _extract_xlsx_text_docling(self, xlsx_path: Path) -> str:
        """Extract text from XLSX using Docling (structured Markdown)."""
        try:
            from docling.document_converter import DocumentConverter

            converter = DocumentConverter()
            result = converter.convert(str(xlsx_path))
            return result.document.export_to_markdown()
        except ImportError:
            logger.info("Docling not installed for XLSX, falling back to openpyxl")
            return ""
        except Exception as e:
            logger.warning("Docling failed for %s: %s, falling back", xlsx_path.name, e)
            return ""

    def _extract_xlsx_text_fallback(self, xlsx_path: Path) -> str:
        """Extract text from XLSX using openpyxl (proper Markdown table syntax)."""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(xlsx_path, read_only=True, data_only=True)
            text_parts: list[str] = []
            for sheet in wb.worksheets:
                sheet_data: list[str] = []
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_text):
                        sheet_data.append("| " + " | ".join(row_text) + " |")
                if sheet_data:
                    ncols = len(sheet_data[0].split(" | "))
                    header_sep = "| " + " | ".join(["---"] * ncols) + " |"
                    sheet_data.insert(1, header_sep)
                    text_parts.append(f"## {sheet.title}\n" + "\n".join(sheet_data))
            return "\n\n".join(text_parts) if text_parts else ""
        except ImportError:
            logger.warning("openpyxl not installed -- XLSX extraction unavailable")
            return ""
        except Exception as e:
            logger.warning("XLSX extraction failed for %s: %s", xlsx_path.name, e)
            return ""

    def _extract_pdf_text(self, pdf_path: Path) -> str:
        """Extract text using pypdf / PyPDF2."""
        try:
            try:
                from pypdf import PdfReader
            except ImportError:
                from PyPDF2 import PdfReader  # type: ignore[no-redef]

            reader = PdfReader(str(pdf_path))
            parts = []
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    parts.append(t)
            return "\n\n".join(parts)
        except Exception as e:
            logger.warning("PDF text extraction failed for %s: %s", pdf_path.name, e)
            return ""

    def _ocr_pdf(self, pdf_path: Path) -> str:
        """OCR a scanned PDF by rendering pages to images, then OCR-ing."""
        try:
            import fitz  # pymupdf

            doc = fitz.open(str(pdf_path))
            parts = []
            for page_num in range(len(doc)):
                pix = doc[page_num].get_pixmap(dpi=300)
                import io

                from PIL import Image

                img = Image.open(io.BytesIO(pix.tobytes("png")))
                parts.append(self._ocr_image_obj(img))
            doc.close()
            return "\n\n".join(p for p in parts if p)
        except ImportError:
            logger.warning("pymupdf or Pillow not installed - OCR fallback unavailable")
            return ""
        except Exception as e:
            logger.warning("OCR fallback failed for %s: %s", pdf_path.name, e)
            return ""

    def _ocr_image(self, image_path: Path) -> str:
        """OCR a single image file."""
        try:
            from PIL import Image

            img = Image.open(image_path)
            return self._ocr_image_obj(img)
        except ImportError:
            logger.warning("Pillow not installed - image OCR unavailable")
            return ""
        except Exception as e:
            logger.warning("Image OCR failed for %s: %s", image_path.name, e)
            return ""

    def _ocr_image_obj(self, img: Any) -> str:
        """Run Tesseract on a PIL Image object."""
        try:
            import pytesseract

            if self.tesseract_path:
                pytesseract.pytesseract.tesseract_cmd = self.tesseract_path
            text: str = pytesseract.image_to_string(img, lang=self.ocr_language)
            return text.strip()
        except ImportError:
            logger.warning("pytesseract not installed")
            return ""
        except Exception as e:
            logger.warning("Tesseract OCR failed: %s", e)
            return ""
